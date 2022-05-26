from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement


def AddGrayscale(shape):

    pic = shape._pic
    rId = pic.xpath('./p:blipFill/a:blip')[0]
    jj = OxmlElement("a:grayscl")
    rId.insert(0,jj)


if __name__ == "__main__":


    file = "grayscale.pptx"


    presentation = Presentation(file)


    startslide = 1
    endslide = len(presentation.slides)

    for sl in range(startslide, endslide + 1):
        slide = presentation.slides[sl - 1]
        for shape in slide.shapes:
            if shape.shape_type == 13:
                AddGrayscale(shape)

    presentation.save(file)
